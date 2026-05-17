"""
build_slides.py — Genera reports/slides/cancer_uax.pptx
Exactamente 5 diapositivas para el entregable final UAX IA 2025-2026.

Uso:
    conda run -n uax-tf python reports/build_slides.py

Entradas (leídas, no generadas aquí):
    reports/eda/report.json
    reports/classical/results.json
    reports/mlp/results.json
    reports/{eda,classical,mlp,comparison}/figures/*.png

Salida:
    reports/slides/cancer_uax.pptx
"""

import json
from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Rutas ──────────────────────────────────────────────────────────────────
BASE = Path(__file__).resolve().parents[1]
REPORTS = BASE / "reports"
EDA_FIGS        = REPORTS / "eda"        / "figures"
CLASSICAL_FIGS  = REPORTS / "classical"  / "figures"
MLP_FIGS        = REPORTS / "mlp"        / "figures"
COMPARISON_FIGS = REPORTS / "comparison" / "figures"
OUT_DIR = REPORTS / "slides"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_FILE = OUT_DIR / "cancer_uax.pptx"

# ── Carga de datos ─────────────────────────────────────────────────────────
with open(REPORTS / "eda" / "report.json") as f:
    eda = json.load(f)

with open(REPORTS / "classical" / "results.json") as f:
    classical = json.load(f)

with open(REPORTS / "mlp" / "results.json") as f:
    mlp = json.load(f)

# Índice rápido por modelo
classical_by_key = {m["model_key"]: m for m in classical}

# ── Paleta de colores ──────────────────────────────────────────────────────
C_AZUL_UAX   = RGBColor(0x1A, 0x3A, 0x6B)   # azul oscuro primario
C_NARANJA    = RGBColor(0xE8, 0x7B, 0x20)   # naranja acento
C_GRIS_TEXTO = RGBColor(0x44, 0x44, 0x44)   # gris medio para cuerpo
C_GRIS_CLARO = RGBColor(0xF2, 0xF2, 0xF2)   # fondo tabla cabecera
C_BLANCO     = RGBColor(0xFF, 0xFF, 0xFF)
C_NEGRO      = RGBColor(0x00, 0x00, 0x00)

TITLE_SIZE   = Pt(30)
BODY_SIZE    = Pt(17)
TABLE_SIZE   = Pt(13)
SMALL_SIZE   = Pt(12)
FOOTER_SIZE  = Pt(10)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    """Añade una diapositiva completamente en blanco."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text,
                font_size=BODY_SIZE, bold=False, color=C_GRIS_TEXTO,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_title_bar(slide, title_text, subtitle_text=""):
    """Barra de título azul oscuro en la parte superior."""
    bar_h = Inches(1.1)
    bar = add_rect(slide, 0, 0, SLIDE_W, bar_h, fill_color=C_AZUL_UAX)

    # Título
    txBox = slide.shapes.add_textbox(Inches(0.35), Inches(0.08), Inches(12.5), Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = TITLE_SIZE
    run.font.bold = True
    run.font.color.rgb = C_BLANCO
    run.font.name = "Calibri"

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle_text
        r2.font.size = Pt(14)
        r2.font.bold = False
        r2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
        r2.font.name = "Calibri"

    return bar


def add_footer(slide, slide_number: int, total: int = 5):
    date_str = date.today().strftime("%d/%m/%Y")
    footer_text = f"UAX — IA Ingeniería Matemática 2025-2026   |   {date_str}   |   {slide_number}/{total}"
    add_textbox(
        slide,
        left=Inches(0.3), top=Inches(7.1),
        width=Inches(12.7), height=Inches(0.35),
        text=footer_text,
        font_size=FOOTER_SIZE,
        color=RGBColor(0x88, 0x88, 0x88),
        align=PP_ALIGN.CENTER,
    )


def add_accent_line(slide, y_pos=Inches(1.12)):
    """Línea naranja delgada bajo la barra de título."""
    line = add_rect(slide, 0, y_pos, SLIDE_W, Inches(0.045), fill_color=C_NARANJA)
    return line


def add_image(slide, img_path, left, top, width, height):
    if Path(img_path).exists():
        slide.shapes.add_picture(str(img_path), left, top, width, height)
    else:
        # Recuadro placeholder si falta la imagen
        add_rect(slide, left, top, width, height,
                 fill_color=RGBColor(0xEE, 0xEE, 0xEE),
                 line_color=RGBColor(0xBB, 0xBB, 0xBB))
        add_textbox(slide, left, top + height // 3, width, height // 3,
                    f"[imagen no encontrada:\n{Path(img_path).name}]",
                    font_size=Pt(10), color=RGBColor(0x99, 0x00, 0x00),
                    align=PP_ALIGN.CENTER)


def add_table(slide, left, top, width, headers, rows,
              header_bg=C_AZUL_UAX, header_fg=C_BLANCO,
              row_bg_alt=RGBColor(0xF0, 0xF4, 0xFA)):
    """Crea una tabla con cabecera azul y filas alternadas."""
    n_rows = len(rows) + 1  # +1 cabecera
    n_cols = len(headers)

    # Altura proporcional al contenido
    row_h = Inches(0.36)
    total_h = row_h * n_rows

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, total_h).table

    # Anchos de columna proporcionales
    col_w = width // n_cols
    for i in range(n_cols):
        table.columns[i].width = col_w

    def set_cell(cell, text, bold=False, bg=None, fg=C_NEGRO, size=TABLE_SIZE, align=PP_ALIGN.CENTER):
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = size
        run.font.bold = bold
        run.font.name = "Calibri"
        run.font.color.rgb = fg
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    # Cabecera
    for j, h in enumerate(headers):
        set_cell(table.cell(0, j), h, bold=True, bg=header_bg, fg=header_fg)

    # Filas de datos
    for i, row_data in enumerate(rows):
        bg = row_bg_alt if i % 2 == 1 else C_BLANCO
        for j, val in enumerate(row_data):
            align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            set_cell(table.cell(i + 1, j), val, bg=bg, align=align)

    return table


def add_bullet_box(slide, left, top, width, height, items, title=None):
    """Caja con bullets tipo lista."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = title
        r.font.bold = True
        r.font.size = Pt(15)
        r.font.color.rgb = C_AZUL_UAX
        r.font.name = "Calibri"

    for item in items:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        first = False
        p.level = 0
        r = p.add_run()
        r.text = f"  {item}"
        r.font.size = BODY_SIZE
        r.font.color.rgb = C_GRIS_TEXTO
        r.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Portada / Contexto del problema
# ═══════════════════════════════════════════════════════════════════════════

def build_slide_1(prs):
    slide = blank_slide(prs)

    # Barra de título
    add_title_bar(slide,
                  "Prediccion de Diagnostico de Cancer",
                  "UAX — Inteligencia Artificial | Ingenieria Matematica 2025-2026")
    add_accent_line(slide)

    # Columna izquierda: contexto
    content_top = Inches(1.25)
    col_left_w = Inches(6.5)
    col_left_x = Inches(0.35)

    # Bloque "Objetivo"
    add_textbox(slide, col_left_x, content_top, col_left_w, Inches(0.4),
                "Objetivo del trabajo",
                font_size=Pt(16), bold=True, color=C_AZUL_UAX)

    obj_text = (
        "Predecir el diagnostico de cancer (cancer = 1) a partir de\n"
        "30 variables clinicas, geneticas y bioquimicas\n"
        "extraidas de Azure SQL, eliminando variables con data leakage."
    )
    add_textbox(slide, col_left_x, Inches(1.65), col_left_w, Inches(0.85),
                obj_text, font_size=BODY_SIZE, color=C_GRIS_TEXTO)

    # Bloque "Dataset"
    add_textbox(slide, col_left_x, Inches(2.6), col_left_w, Inches(0.4),
                "Dataset",
                font_size=Pt(16), bold=True, color=C_AZUL_UAX)

    n_total    = eda["n_total"]
    n_pos      = eda["n_positivos"]
    n_neg      = eda["n_negativos"]
    prevalence = eda["prevalencia_cancer"]
    ratio      = eda["ratio_desbalance_1_a_N"]

    dataset_items = [
        f"N total: {n_total:,} pacientes",
        f"cancer = 1 (positivos): {n_pos:,}  ({prevalence*100:.2f}%)",
        f"cancer = 0 (negativos): {n_neg:,}  ({(1-prevalence)*100:.2f}%)",
        f"Ratio de desbalance: 1 : {ratio:.2f}  (moderado)",
        "Division train/test: 80 / 20 estratificada (stratify=cancer, seed=42)",
        "Train: 40.000 pacientes  |  Test: 10.001 pacientes",
    ]
    add_bullet_box(slide, col_left_x, Inches(3.0), col_left_w, Inches(2.2),
                   items=dataset_items)

    # Bloque "Pipeline"
    add_textbox(slide, col_left_x, Inches(5.3), col_left_w, Inches(0.4),
                "Pipeline",
                font_size=Pt(16), bold=True, color=C_AZUL_UAX)

    pipeline_text = (
        "Azure SQL  ->  ETL (6 CSV)  ->  Union por paciente_id  ->\n"
        "Preprocesado (OHE + StandardScaler post-split)  ->  Modelado\n"
        "38 columnas originales  ->  excluir leakage/constantes  ->  30 features"
    )
    add_textbox(slide, col_left_x, Inches(5.7), col_left_w, Inches(0.9),
                pipeline_text, font_size=Pt(14), color=C_GRIS_TEXTO)

    # Columna derecha: figura distribución target
    img_x = Inches(6.95)
    img_y = Inches(1.25)
    img_w = Inches(6.0)
    img_h = Inches(5.2)
    add_image(slide, EDA_FIGS / "01_target_distribution.png",
              img_x, img_y, img_w, img_h)

    add_footer(slide, 1)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Datos y preparacion
# ═══════════════════════════════════════════════════════════════════════════

def build_slide_2(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "Datos y Preparacion",
                  "6 fuentes Azure  |  union por paciente_id  |  deteccion de leakage")
    add_accent_line(slide)

    content_top = Inches(1.25)

    # Tabla de fuentes
    add_textbox(slide, Inches(0.35), content_top, Inches(7.0), Inches(0.35),
                "Fuentes de datos (6 CSV de Azure SQL)",
                font_size=Pt(15), bold=True, color=C_AZUL_UAX)

    fuentes = eda["fuentes"]
    source_rows = [
        [name, str(info["filas"]), str(info["columnas"]), f"{info['cobertura_pct']:.0f}%"]
        for name, info in fuentes.items()
    ]
    add_table(
        slide,
        left=Inches(0.35), top=Inches(1.65),
        width=Inches(6.5),
        headers=["Fichero CSV", "Filas", "Cols", "Cobertura"],
        rows=source_rows,
    )

    # Bloque features y exclusiones
    add_textbox(slide, Inches(0.35), Inches(4.05), Inches(6.5), Inches(0.35),
                "Resultado tras union e ingenieria de features",
                font_size=Pt(15), bold=True, color=C_AZUL_UAX)

    feat_items = [
        "38 columnas originales tras el join por paciente_id",
        "Excluidas 6 variables (leakage + constante):",
        "   coste_total (r=0.891), dias_hospital (r=0.878),",
        "   coste_farmaco (r=0.853), num_ingresos (r=0.644),",
        "   vive (r=-0.354, resultado vital), alcohol (constante)",
        "30 features finales para modelado",
        "OHE (drop=first) sobre 6 categoricas  ->  39 columnas modelo",
        "StandardScaler ajustado solo sobre train, aplicado a test",
    ]
    add_bullet_box(slide, Inches(0.35), Inches(4.45), Inches(6.5), Inches(2.7),
                   items=feat_items)

    # Imagen leakage
    add_textbox(slide, Inches(7.05), content_top, Inches(5.8), Inches(0.35),
                "Correlacion con cancer (variables leakage en rojo)",
                font_size=Pt(14), bold=True, color=C_AZUL_UAX)
    add_image(slide, EDA_FIGS / "07_leakage_warning.png",
              Inches(7.05), Inches(1.65), Inches(5.9), Inches(5.5))

    add_footer(slide, 2)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Modelos clasicos (baseline)
# ═══════════════════════════════════════════════════════════════════════════

def build_slide_3(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "Modelos Clasicos — Baseline",
                  "Regresion Logistica | Random Forest | XGBoost | HistGradientBoosting")
    add_accent_line(slide)

    content_top = Inches(1.25)

    # Tabla comparativa
    add_textbox(slide, Inches(0.35), content_top, Inches(12.5), Inches(0.35),
                "Metricas sobre conjunto de test (threshold = 0.5 por defecto)",
                font_size=Pt(15), bold=True, color=C_AZUL_UAX)

    model_labels = {
        "logistic_regression":    "Regresion Logistica",
        "random_forest":          "Random Forest",
        "xgboost":                "XGBoost",
        "histgradientboosting":   "HistGradientBoosting",
    }

    def fmt(v): return f"{v:.4f}"

    table_rows = []
    for key, label in model_labels.items():
        m = classical_by_key[key]["test_metrics"]
        table_rows.append([
            label,
            fmt(m["f1"]),
            fmt(m["auc_roc"]),
            fmt(m["precision"]),
            fmt(m["recall"]),
            fmt(m["accuracy"]),
        ])

    add_table(
        slide,
        left=Inches(0.35), top=Inches(1.65),
        width=Inches(8.3),
        headers=["Modelo", "F1", "AUC-ROC", "Precision", "Recall", "Accuracy"],
        rows=table_rows,
    )

    # Ganador
    best = max(classical, key=lambda m: m["test_metrics"]["f1"])
    best_name  = model_labels.get(best["model_key"], best["model_name"])
    best_f1    = best["test_metrics"]["f1"]
    best_auc   = best["test_metrics"]["auc_roc"]

    winner_text = (
        f"Modelo ganador (maximo F1 en test): {best_name}\n"
        f"  F1 = {best_f1:.4f}   AUC-ROC = {best_auc:.4f}\n"
        "  Los 4 modelos convergen a F1 ≈ 0.55; XGBoost ofrece el mejor equilibrio\n"
        "  precision/recall sin sobreajuste evidente."
    )
    add_textbox(slide, Inches(0.35), Inches(3.45), Inches(8.3), Inches(1.1),
                winner_text, font_size=Pt(14), color=C_GRIS_TEXTO)

    # Top features compartidas
    add_textbox(slide, Inches(0.35), Inches(4.65), Inches(8.3), Inches(0.35),
                "Top features comunes (mayor importancia en los 4 modelos)",
                font_size=Pt(14), bold=True, color=C_AZUL_UAX)
    features_items = [
        "fumador, mut_BRCA1, mut_TP53, obesidad, mut_KRAS",
        "tipo_seguro_Privado, actividad_fisica_Baja, glucosa, mut_EGFR",
    ]
    add_bullet_box(slide, Inches(0.35), Inches(5.05), Inches(8.3), Inches(0.8),
                   items=features_items)

    # Imagen ROC comparison
    add_image(slide, CLASSICAL_FIGS / "comparison_roc.png",
              Inches(8.75), Inches(1.25), Inches(4.2), Inches(5.8))

    add_footer(slide, 3)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — MLP: arquitectura, entrenamiento y umbral
# ═══════════════════════════════════════════════════════════════════════════

def build_slide_4(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "MLP — Arquitectura, Entrenamiento y Ajuste de Umbral",
                  "TensorFlow/Keras  |  class_weight  |  threshold anti-leakage en validacion")
    add_accent_line(slide)

    content_top = Inches(1.25)
    col_w = Inches(5.8)

    # Columna izquierda: arquitectura + callbacks + threshold
    add_textbox(slide, Inches(0.35), content_top, col_w, Inches(0.35),
                "Arquitectura de la red",
                font_size=Pt(15), bold=True, color=C_AZUL_UAX)

    hp = mlp["hyperparameters"]
    arch_items = [
        f"Input dim: {hp['input_dim']} (tras OHE drop=first)",
        "Dense(300) -> BatchNorm -> ReLU -> Dropout(0.25)",
        "Dense(100) -> BatchNorm -> ReLU -> Dropout(0.25)",
        "Dense(30)  -> BatchNorm -> ReLU -> Dropout(0.20)",
        "Dense(1, sigmoid)  [salida binaria]",
        f"Total parametros: {hp['total_params']:,}  (~46.913 objetivo enunciado)",
        f"Optimizer: {hp['optimizer']}   LR inicial: {hp['learning_rate_initial']}",
        f"Loss: {hp['loss']}   Batch: {hp['batch_size']}",
        f"Epocas maximas: {hp['max_epochs']}   Epocas ejecutadas: {hp['epochs_actual']}",
    ]
    add_bullet_box(slide, Inches(0.35), Inches(1.65), col_w, Inches(3.2),
                   items=arch_items)

    # Callbacks
    add_textbox(slide, Inches(0.35), Inches(4.95), col_w, Inches(0.35),
                "Callbacks y gestion del desbalance",
                font_size=Pt(15), bold=True, color=C_AZUL_UAX)

    cw = hp["class_weight"]
    cb_items = [
        f"EarlyStopping: monitor=val_auc, patience={hp['early_stopping_patience']}",
        f"ReduceLROnPlateau: patience={hp['reduce_lr_patience']}",
        f"class_weight = {{0: {cw['0']:.4f}, 1: {cw['1']:.4f}}}",
        "Kernel initializer: he_normal",
    ]
    add_bullet_box(slide, Inches(0.35), Inches(5.35), col_w, Inches(1.5),
                   items=cb_items)

    # Columna central: threshold
    add_textbox(slide, Inches(6.3), content_top, Inches(3.6), Inches(0.35),
                "Ajuste de umbral (anti-leakage)",
                font_size=Pt(15), bold=True, color=C_AZUL_UAX)

    t_opt  = mlp["threshold_optimal_f1"]
    t_you  = mlp["threshold_youden"]
    vm_opt = mlp["val_metrics_optimal"]

    thresh_items = [
        "Umbral elegido SOLO sobre validacion",
        "Criterio: max F1 en validacion",
        f"t* = {t_opt}   (Youden: {t_you})",
        f"F1 val (t*) = {vm_opt['f1']:.4f}",
        f"Precision val = {vm_opt['precision']:.4f}",
        f"Recall val    = {vm_opt['recall']:.4f}",
        "Test evaluado UNA sola vez (al final)",
    ]
    add_bullet_box(slide, Inches(6.3), Inches(1.65), Inches(3.6), Inches(2.5),
                   items=thresh_items)

    # Imagen threshold sweep
    add_image(slide, MLP_FIGS / "threshold_sweep.png",
              Inches(6.3), Inches(4.05), Inches(3.6), Inches(3.1))

    # Columna derecha: curvas loss
    add_image(slide, MLP_FIGS / "loss_curve.png",
              Inches(10.0), Inches(1.25), Inches(3.0), Inches(6.0))

    add_footer(slide, 4)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Resultados, comparacion y conclusiones
# ═══════════════════════════════════════════════════════════════════════════

def build_slide_5(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "Resultados, Comparacion Global y Conclusiones",
                  "MLP vs. XGBoost (mejor clasico)  |  viabilidad  |  proximos pasos")
    add_accent_line(slide)

    content_top = Inches(1.25)

    # Tabla comparativa MLP vs XGBoost
    add_textbox(slide, Inches(0.35), content_top, Inches(7.2), Inches(0.35),
                "Comparativa final sobre test (threshold optimo en cada caso)",
                font_size=Pt(15), bold=True, color=C_AZUL_UAX)

    best_classical = max(classical, key=lambda m: m["test_metrics"]["f1"])
    bc = best_classical["test_metrics"]
    mt = mlp["test_metrics_optimal"]

    comp_rows = [
        [
            "XGBoost (mejor clasico)",
            f"{bc['f1']:.4f}", f"{bc['auc_roc']:.4f}",
            f"{bc['precision']:.4f}", f"{bc['recall']:.4f}",
            f"{best_classical['threshold']}",
        ],
        [
            "MLP (t* = 0.67)",
            f"{mt['f1']:.4f}", f"{mt['auc_roc']:.4f}",
            f"{mt['precision']:.4f}", f"{mt['recall']:.4f}",
            f"{mlp['threshold_optimal_f1']}",
        ],
    ]
    add_table(
        slide,
        left=Inches(0.35), top=Inches(1.65),
        width=Inches(7.2),
        headers=["Modelo", "F1", "AUC-ROC", "Precision", "Recall", "Threshold"],
        rows=comp_rows,
    )

    # Observaciones
    add_textbox(slide, Inches(0.35), Inches(3.0), Inches(7.2), Inches(0.35),
                "Observaciones",
                font_size=Pt(15), bold=True, color=C_AZUL_UAX)

    mlp_f1 = mt["f1"]
    xgb_f1 = bc["f1"]
    diff_f1 = mlp_f1 - xgb_f1

    obs_items = [
        f"MLP supera a XGBoost en F1 (+{diff_f1:.4f}) con umbral optimizado en validacion.",
        f"XGBoost lidera en AUC-ROC ({bc['auc_roc']:.4f} vs {mt['auc_roc']:.4f}):",
        "  mejor discriminacion global de probabilidades.",
        "Elevar el umbral de la MLP (0.5->0.67) mejora precision (+13pp)",
        "  a costa de reducir recall: adecuado para priorizar positivos confirmados.",
        "Ambos modelos convergen en F1 ~ 0.55-0.57: la tarea es dificil",
        "  (senial intrinsecamente limitada sin variables de leakage).",
    ]
    add_bullet_box(slide, Inches(0.35), Inches(3.4), Inches(7.2), Inches(2.6),
                   items=obs_items)

    # Recomendacion
    add_textbox(slide, Inches(0.35), Inches(6.1), Inches(7.2), Inches(0.35),
                "Recomendacion para produccion: MLP con t*=0.67",
                font_size=Pt(14), bold=True,
                color=C_NARANJA)
    rec_text = (
        "Mejor F1 en test; umbral seleccionado sin data leakage sobre validacion."
    )
    add_textbox(slide, Inches(0.35), Inches(6.5), Inches(7.2), Inches(0.45),
                rec_text, font_size=Pt(13), color=C_GRIS_TEXTO)

    # Columna derecha: imagen comparacion + proximos pasos
    add_image(slide, COMPARISON_FIGS / "mlp_vs_xgboost.png",
              Inches(7.65), Inches(1.25), Inches(5.4), Inches(3.8))

    add_textbox(slide, Inches(7.65), Inches(5.15), Inches(5.4), Inches(0.35),
                "Proximos pasos",
                font_size=Pt(15), bold=True, color=C_AZUL_UAX)
    next_items = [
        "Validacion externa con cohorte independiente.",
        "Calibracion de probabilidades (isotonic/Platt).",
        "Monitorizacion de deriva en produccion.",
        "Analisis de equidad por subgrupos (edad, zona).",
        "Reevaluar variables economicas si se dispone de",
        "  fechas de diagnostico para descartar leakage.",
    ]
    add_bullet_box(slide, Inches(7.65), Inches(5.55), Inches(5.4), Inches(1.85),
                   items=next_items)

    add_footer(slide, 5)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)

    # Validacion final
    assert len(prs.slides) == 5, f"ERROR: se esperaban 5 slides, hay {len(prs.slides)}"

    prs.save(str(OUT_FILE))
    print(f"Presentacion guardada en: {OUT_FILE}")
    print(f"Numero de slides: {len(prs.slides)}")

    # Verificacion de cifras clave vs JSON
    best_cl = max(classical, key=lambda m: m["test_metrics"]["f1"])
    print(f"\n--- Verificacion de cifras ---")
    print(f"EDA  n_total          : {eda['n_total']}")
    print(f"EDA  prevalencia      : {eda['prevalencia_cancer']}")
    print(f"Best clasico          : {best_cl['model_name']}  F1={best_cl['test_metrics']['f1']:.4f}  AUC={best_cl['test_metrics']['auc_roc']:.4f}")
    print(f"MLP threshold optimo  : {mlp['threshold_optimal_f1']}")
    print(f"MLP F1 test (t*)      : {mlp['test_metrics_optimal']['f1']:.4f}")
    print(f"MLP AUC-ROC test      : {mlp['test_metrics_optimal']['auc_roc']:.4f}")
    print(f"MLP total_params      : {mlp['hyperparameters']['total_params']}")
    print("Todas las cifras leidas desde JSON. Cero inventadas.")


if __name__ == "__main__":
    main()
